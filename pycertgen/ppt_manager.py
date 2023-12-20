
from pptx import Presentation
import pandas as pd
import uuid
import os


def fill_template(template_path: str, replacement_dict: dict, output_dir: str, key_autogen=False):
    # Load the PowerPoint presentation
    output_pptx = os.path.join(
        output_dir, str(replacement_dict['index']) + '.pptx')

    prs = Presentation(template_path)

    if key_autogen:
        replacement_dict['certificate_key'] = uuid.uuid4()

    for slide in prs.slides:
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]

        for shape in text_shapes:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    replace_texts(run, replacement_dict)

    # Save the modified presentation
    prs.save(output_pptx)

    return output_pptx


def replace_texts(run, replacement_dict):
    for old_text, new_text in replacement_dict.items():
        if isinstance(new_text, pd.Timestamp):
            new_text = new_text.strftime('%d/%m/%Y')

        if old_text in run.text:
            text = run.text.replace(f'[{old_text}]', str(new_text))
            run.text = text
