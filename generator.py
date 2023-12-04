from pptx import Presentation
import pandas as pd
import uuid
import os
import subprocess
from typing import Literal


def replace_text_in_ppt(template_path, replacement_dict, key_autogen=False,
                        to_pdf_method: Literal['libreoffice', 'powerpoint'] = None):
    # Load the PowerPoint presentation
    output_pptx = os.path.join(
        'output', str(replacement_dict['index']) + '.pptx')

    output_pdf = os.path.join(
        'output', str(replacement_dict['index']) + '.pdf')

    prs = Presentation(template_path)

    if key_autogen:
        replacement_dict['certificate_key'] = uuid.uuid4()

    for slide in prs.slides:
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]

        for shape in text_shapes:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    replace_texts(run, replacement_dict)

    # Save the modified presentation
    prs.save(output_pptx)

    if to_pdf_method:
        convert_ppt_to_pdf(output_pptx, output_pdf, to_pdf_method)


def replace_texts(run, replacement_dict):
    for old_text, new_text in replacement_dict.items():
        if isinstance(new_text, pd.Timestamp):
            new_text = new_text.strftime('%d/%m/%Y')

        if old_text in run.text:
            text = run.text.replace(f'[{old_text}]', str(new_text))
            run.text = text


def convert_ppt_to_pdf(input_ppt, output_dir, method: Literal['libreoffice', 'powerpoint'] = 'libreoffice'):
    try:
        if method == 'libreoffice':
            subprocess.run(['libreoffice', '--headless', '--invisible',
                            '--convert-to', 'pdf', '--outdir', os.path.dirname(output_dir), input_ppt])

        elif method == 'powerpoint':
            raise NotImplementedError("powerpoint method not implemented yet.")

    except subprocess.CalledProcessError as e:
        print(f"Conversion failed: {e}")


def gen(certified_list_path: str, template_path: str, key_autogen=False, to_pdf_method: Literal['libreoffice', 'powerpoint'] = None):
    certified_df = pd.read_excel(certified_list_path)

    certified_data = certified_df.reset_index().to_dict(orient='records')

    for certified in certified_data:

        # Call the function to replace text in the PowerPoint presentation
        replace_text_in_ppt(template_path, certified,
                            key_autogen, to_pdf_method)
